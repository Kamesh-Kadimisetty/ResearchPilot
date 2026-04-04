import re
from pypdf import PdfReader
from pptx import Presentation
from io import BytesIO

def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from a PDF file."""
    try:
        reader = PdfReader(BytesIO(file_content))
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        return f"Error extracting PDF: {str(e)}"

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from a PPTX file."""
    try:
        prs = Presentation(BytesIO(file_content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        return f"Error extracting PPTX: {str(e)}"

def clean_input(text: str) -> str:
    """Clean and normalize input text."""
    if not text:
        return ""
    # Remove excessive whitespace
    text = re.sub(r'\s+', ' ', text)
    # Basic sanitization
    return text.strip()

def latex_escape(text: str) -> str:
    """
    Escape special characters for LaTeX.
    """
    if not text:
        return ""
    
    # Map of special characters to their LaTeX escape sequences
    mapping = {
        '&': r'\&',
        '%': r'\%',
        '$': r'\$',
        '#': r'\#',
        '_': r'\_',
        '{': r'\{',
        '}': r'\}',
        '~': r'\textasciitilde{}',
        '^': r'\textasciicircum{}',
        '\\': r'\textbackslash{}',
    }
    
    # regex pattern to match any of the keys in mapping
    regex = re.compile('|'.join(re.escape(str(key)) for key in sorted(mapping.keys(), key=lambda item: -len(item))))
    
    return regex.sub(lambda match: mapping[match.group()], text)
